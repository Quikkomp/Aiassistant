# document_loader.py
import os
from docx import Document
from PyPDF2 import PdfReader, errors
from pptx import Presentation


def read_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def read_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        if not text.strip():
            print(f"⚠️ PDF {file_path} 没有可提取的文本")
        return text
    except errors.PdfReadError:
        print(f"❌ 无法读取 PDF 文件: {file_path}")
        return ""
    except errors.EmptyFileError:
        print(f"❌ PDF 文件为空: {file_path}")
        return ""

def read_docx(file_path):
    try:
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        print(f"❌ 无法读取 DOCX 文件 {file_path}: {e}")
        return ""

def read_pptx(file_path):
    """
    读取 PPTX 中所有文本框、标题等的文字，按换行拼成一个长字符串。
    """
    try:
        prs = Presentation(file_path)
        texts = []

        # 遍历所有幻灯片 & 形状
        for slide in prs.slides:
            for shape in slide.shapes:
                # 大部分带文字的 shape 都有 .text 属性
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)

        # 也可以根据需要再补充读取 notes、表格等
        if not texts:
            print(f"⚠️ PPTX {file_path} 中没有提取到任何文本")

        return "\n".join(texts)

    except Exception as e:
        print(f"❌ 无法读取 PPTX 文件 {file_path}: {e}")
        return ""

def load_documents(folder_path="materials"):
    all_docs = {}

    if not os.path.exists(folder_path):
        print(f"⚠️ 路径不存在：{folder_path}")
        return all_docs

    for name in os.listdir(folder_path):
        path = os.path.join(folder_path, name)

        # ⭐ 多用户模式下 materials 里会有 user_xxxx 子文件夹，这里要跳过目录
        if os.path.isdir(path):
            # 如果你希望递归子文件夹，这里可以改成 os.walk；目前我们只加载这一层
            continue

        lower = name.lower()
        if lower.endswith(".txt"):
            all_docs[name] = read_txt(path)
        elif lower.endswith(".pdf"):
            all_docs[name] = read_pdf(path)
        elif lower.endswith(".docx"):
            all_docs[name] = read_docx(path)
        elif lower.endswith(".pptx"):
            all_docs[name] = read_pptx(path)  # ← 新增这行
        else:
            print(f"⚠️ 不支持的文件格式：{name}")

    return all_docs

